"""
Rebuild BM project-human data_update.pptx with actual figures embedded.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

FIG_DIR = "bm_analysis_out/figures"
OUT_PPTX = "BM project-human data_update.pptx"

TITLE_COLOR = RGBColor(0x1F, 0x49, 0x7D)
BODY_COLOR  = RGBColor(0x26, 0x26, 0x26)
GRAY        = RGBColor(0x88, 0x88, 0x88)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = TITLE_COLOR; bar.line.fill.background()
    txb = slide.shapes.add_textbox(Inches(0.3), Inches(0.12), Inches(12.0), Inches(0.55))
    p = txb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
    if subtitle:
        txb2 = slide.shapes.add_textbox(Inches(0.3), Inches(0.68), Inches(12.0), Inches(0.35))
        p2 = txb2.text_frame.paragraphs[0]
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(13); r2.font.color.rgb = RGBColor(0xCC,0xDD,0xEE)

def add_img(slide, path, left, top, width, height):
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)
    else:
        box = slide.shapes.add_shape(1, left, top, width, height)
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF0,0xF4,0xF8)
        box.line.color.rgb = RGBColor(0xAA,0xBB,0xCC)

def body_text(slide, left, top, width, height, lines, font_size=13):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run(); r.text = line
        if line.startswith("•") or line.startswith("  –"):
            r.font.size = Pt(font_size); r.font.color.rgb = BODY_COLOR
        else:
            r.font.size = Pt(font_size); r.font.bold = True; r.font.color.rgb = TITLE_COLOR

# ── Slide 1: Title ────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
bg = slide.shapes.add_shape(1,0,0,Inches(13.33),Inches(7.5))
bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0xF7,0xF9,0xFC); bg.line.fill.background()
bar = slide.shapes.add_shape(1,0,Inches(2.8),Inches(13.33),Inches(2.0))
bar.fill.solid(); bar.fill.fore_color.rgb = TITLE_COLOR; bar.line.fill.background()
txb = slide.shapes.add_textbox(Inches(1.0),Inches(3.0),Inches(11.0),Inches(1.2))
p = txb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "BM Project – Human Data Update"
r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
txb2 = slide.shapes.add_textbox(Inches(1.0),Inches(4.1),Inches(11.0),Inches(0.5))
p2 = txb2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = "Macrophage Sub-clustering Analysis | GSE266330 | 30 samples | 15,074 macrophages"
r2.font.size = Pt(15); r2.font.color.rgb = RGBColor(0xCC,0xDD,0xEE)
txb3 = slide.shapes.add_textbox(Inches(1.0),Inches(6.5),Inches(11.0),Inches(0.5))
p3 = txb3.text_frame.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run(); r3.text = "Jianting Sheng | Houston Methodist Neal Cancer Center"
r3.font.size = Pt(13); r3.font.color.rgb = GRAY

# ── Slide 2: UMAP overview ────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
title_bar(slide, "Macrophage Sub-clustering (UMAP)",
          "30 samples pooled → 13 sub-clusters | GSE266330, Liu et al. 2025")
add_img(slide, os.path.join(FIG_DIR, "01_umap_overview.png"),
        Inches(3.0), Inches(1.2), Inches(7.3), Inches(6.1))

# ── Slide 3: Cancer type composition ─────────────────────────────────────────
slide = prs.slides.add_slide(blank)
title_bar(slide, "Cancer Type Composition per Macrophage Cluster",
          "BC / BDC / BM / CC / EC / KC samples")
add_img(slide, os.path.join(FIG_DIR, "02_cancer_type_bar.png"),
        Inches(1.0), Inches(1.2), Inches(11.0), Inches(6.0))

# ── Slide 4: ESR1 + SP1 ───────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
title_bar(slide, "ER / SP1 Co-expression in Macrophage Sub-clusters",
          "Cluster 9 identified as ER+/SP1-high")
add_img(slide, os.path.join(FIG_DIR, "03a_feature_plots_ESR1_SP1.png"),
        Inches(0.3), Inches(1.2), Inches(8.5), Inches(4.5))
add_img(slide, os.path.join(FIG_DIR, "03b_dotplot_ESR1_SP1.png"),
        Inches(0.3), Inches(5.8), Inches(8.5), Inches(1.5))
body_text(slide, Inches(9.0), Inches(1.3), Inches(4.0), Inches(5.5), [
    "Key finding",
    "• Cluster 9 = ER+/SP1-high",
    "• Both ESR1 and SP1 expressed",
    "  – above background in cluster 9",
    "• Matches Yuliang's gene signature",
    "",
    "Next",
    "• Validate with marker genes",
    "• Cross-ref S2C2 crosstalk axis",
], font_size=12)

# ── Slide 5: Metastasis vs. Control ──────────────────────────────────────────
slide = prs.slides.add_slide(blank)
title_bar(slide, "Macrophage Cluster Enrichment — Metastasis vs. Control",
          "Are specific sub-clusters over-represented in metastatic samples?")
add_img(slide, os.path.join(FIG_DIR, "04_metastasis_vs_control.png"),
        Inches(0.5), Inches(1.2), Inches(12.3), Inches(6.1))

# ── Slide 6: ER+SP1 co-expression UMAP ───────────────────────────────────────
slide = prs.slides.add_slide(blank)
title_bar(slide, "ER + SP1 Co-expression Score — UMAP & Violin",
          "Combined ESR1+SP1 module score; cluster 9 highest mean")
add_img(slide, os.path.join(FIG_DIR, "05_ER_SP1_coexpression.png"),
        Inches(0.3), Inches(1.2), Inches(12.7), Inches(6.1))

# ── Slide 7: DEG heatmap ─────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
title_bar(slide, "Top DEGs per Macrophage Sub-cluster (Wilcoxon)",
          "Top 4 marker genes per cluster")
add_img(slide, os.path.join(FIG_DIR, "06_deg_heatmap.png"),
        Inches(0.3), Inches(1.2), Inches(12.7), Inches(6.1))

# ── Slide 8: Gender + tissue origin ─────────────────────────────────────────
slide = prs.slides.add_slide(blank)
title_bar(slide, "Gender & Tissue Origin per Macrophage Cluster")
add_img(slide, os.path.join(FIG_DIR, "06_gender_tissue_origin.png"),
        Inches(0.3), Inches(1.2), Inches(12.7), Inches(6.1))

# ── Slide 8: Summary ─────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
title_bar(slide, "Summary & Next Steps")
div = slide.shapes.add_shape(1,Inches(6.7),Inches(1.3),Inches(0.03),Inches(5.5))
div.fill.solid(); div.fill.fore_color.rgb = RGBColor(0xCC,0xDD,0xEE); div.line.fill.background()
body_text(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.5), [
    "Results",
    "• 13 macrophage sub-clusters identified",
    "  – from 15,074 macrophages, 30 samples",
    "• Cluster 9 = ER+/SP1-high",
    "  – highest combined ESR1+SP1 score",
    "• BC samples dominant in cluster 9",
    "  – consistent with ER+ breast cancer biology",
    "• Gender + tissue origin breakdown shown",
], font_size=13)
body_text(slide, Inches(7.0), Inches(1.3), Inches(6.0), Inches(5.5), [
    "Next steps",
    "• Annotate clusters with canonical markers",
    "  – Cross-ref Yuliang's gene signature",
    "• Link cluster 9 to S2C2 crosstalk model",
    "  – ER+ Mφ → CD8+ T cell axis",
    "• Validate in mouse scRNA (Jianting W2/W3)",
    "• Drug target screen on cluster 9",
    "• Obtain remaining 32 samples (LC/PC/RC/TC)",
    "  – pending storage quota expansion",
], font_size=13)

prs.save(OUT_PPTX)
print(f"Saved: {OUT_PPTX}")
