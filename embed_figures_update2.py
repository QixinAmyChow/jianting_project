"""
BM project-human data_update_2.pptx — 1 figure per slide
Figures from bm_analysis_out/figures_update_2/
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

FIG_DIR  = "bm_analysis_out/figures_update_2"
OUT_PPTX = "BM project-human data_update_2.pptx"

_GENE_FILE = "bm_analysis_out/figures/highlight_genes.txt"
try:
    raw = open(_GENE_FILE).read().strip()
    HIGHLIGHT_GENES = [g.strip() for g in raw.split(",") if g.strip()]
except FileNotFoundError:
    HIGHLIGHT_GENES = ["Timp2", "Mmp14", "Pdgfa", "Fn1", "Vim"]
GENE_LABEL = " / ".join(HIGHLIGHT_GENES)

TITLE_COLOR = RGBColor(0x1F, 0x49, 0x7D)
BODY_COLOR  = RGBColor(0x26, 0x26, 0x26)
GRAY        = RGBColor(0x88, 0x88, 0x88)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = TITLE_COLOR; bar.line.fill.background()
    txb = slide.shapes.add_textbox(Inches(0.3), Inches(0.12), Inches(12.0), Inches(0.55))
    p = txb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
    if subtitle:
        txb2 = slide.shapes.add_textbox(Inches(0.3), Inches(0.68), Inches(12.0), Inches(0.35))
        p2 = txb2.text_frame.paragraphs[0]
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(13); r2.font.color.rgb = RGBColor(0xCC,0xDD,0xEE)

def add_img(slide, path, left, top, width, height):
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)
    else:
        box = slide.shapes.add_shape(1, left, top, width, height)
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF0,0xF4,0xF8)
        box.line.color.rgb = RGBColor(0xAA,0xBB,0xCC)
        txb = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.1),
                                        width - Inches(0.2), height - Inches(0.2))
        r = txb.text_frame.paragraphs[0].add_run()
        r.text = f"[Missing: {os.path.basename(path)}]"
        r.font.size = Pt(11); r.font.color.rgb = GRAY

def fig_slide(title, subtitle, img_path):
    slide = prs.slides.add_slide(blank)
    title_bar(slide, title, subtitle)
    add_img(slide, img_path, Inches(0.4), Inches(1.2), Inches(12.5), Inches(6.1))
    return slide

# ── Slide 1: Title ────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
bg = slide.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(7.5))
bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0xF7,0xF9,0xFC); bg.line.fill.background()
bar = slide.shapes.add_shape(1, 0, Inches(2.8), Inches(13.33), Inches(2.0))
bar.fill.solid(); bar.fill.fore_color.rgb = TITLE_COLOR; bar.line.fill.background()
txb = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.0), Inches(1.2))
p = txb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "BM Project \u2013 Human Data Update 2"
r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
txb2 = slide.shapes.add_textbox(Inches(1.0), Inches(4.1), Inches(11.0), Inches(0.5))
p2 = txb2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "Macrophage Sub-clustering | GSE266330 | ER status from patient metadata"
r2.font.size = Pt(15); r2.font.color.rgb = RGBColor(0xCC,0xDD,0xEE)
txb3 = slide.shapes.add_textbox(Inches(1.0), Inches(6.5), Inches(11.0), Inches(0.5))
p3 = txb3.text_frame.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run(); r3.text = "Jianting Sheng | Houston Methodist Neal Cancer Center"
r3.font.size = Pt(13); r3.font.color.rgb = GRAY

# ── Slide 2: UMAP ─────────────────────────────────────────────────────────────
fig_slide("Macrophage Sub-clustering (UMAP)",
          "Full recompute: SCTransform \u2192 PCA \u2192 UMAP \u2192 FindClusters | GSE266330, Liu et al. 2025",
          os.path.join(FIG_DIR, "01_umap_overview.png"))

# ── Slide 3: Cancer type bar ──────────────────────────────────────────────────
fig_slide("Cancer Type Composition per Macrophage Cluster",
          "BC split by ER status from patient metadata (Table S1)",
          os.path.join(FIG_DIR, "02_cancer_type_bar_bc_er.png"))

# ── Slide 4: Feature plots ────────────────────────────────────────────────────
fig_slide(f"Highlight Gene Feature Plots \u2014 UMAP",
          f"Genes: {GENE_LABEL}",
          os.path.join(FIG_DIR, "03a_feature_plots_highlight.png"))

# ── Slide 5: Dot plot (highlight genes) ───────────────────────────────────────
fig_slide(f"Highlight Gene Dot Plot \u2014 per Macrophage Sub-cluster",
          f"Genes: {GENE_LABEL} | Dot size = % expressing; color = avg expression",
          os.path.join(FIG_DIR, "03b_dotplot_highlight.png"))

# ── Slide 6: DEG heatmap ──────────────────────────────────────────────────────
fig_slide("DEG Heatmap \u2014 Top 5 Marker Genes per Cluster",
          "FindAllMarkers (Wilcoxon) | only.pos=TRUE | min.pct=0.25 | padj<0.05",
          os.path.join(FIG_DIR, "04a_deg_heatmap.png"))

# ── Slide 7: DEG dot plot ─────────────────────────────────────────────────────
fig_slide("DEG Dot Plot \u2014 Top 3 Marker Genes per Cluster",
          "Dot size = % cells expressing; color = average expression",
          os.path.join(FIG_DIR, "04b_deg_dotplot.png"))

# ── Slide 8: GSEA NES heatmap ─────────────────────────────────────────────────
fig_slide("GSEA \u2014 Hallmark + Oncogenic (C6) NES Heatmap",
          "fgsea | MSigDB Hallmark (H) + Oncogenic Signatures (C6) | HVGs only | color = NES",
          os.path.join(FIG_DIR, "05a_gsea_nes_heatmap.png"))

# ── Slide 9: GSEA bubble plot ─────────────────────────────────────────────────
fig_slide("GSEA \u2014 Top Pathways per Cluster (Bubble Plot)",
          "Top 3 pathways per cluster (padj<0.25) | size = \u2212log\u2081\u2080(padj) | color = NES",
          os.path.join(FIG_DIR, "05b_gsea_bubble.png"))

# ── Slide 10: Gender per cluster ──────────────────────────────────────────────
fig_slide("Gender Distribution per Macrophage Sub-cluster",
          "Female / Male / Unknown fraction per cluster",
          os.path.join(FIG_DIR, "06_gender_per_cluster.png"))

# ── Slide 11: Tissue origin per cluster ───────────────────────────────────────
fig_slide("Bone Site / Tissue Origin per Macrophage Sub-cluster",
          "Site of bone metastasis resection",
          os.path.join(FIG_DIR, "07_tissue_origin_per_cluster.png"))

# ── Slide 12: Treatment response per cluster ──────────────────────────────────
fig_slide("Treatment Response per Macrophage Sub-cluster",
          "Neoadjuvant or metastatic treatment response (Table S1)",
          os.path.join(FIG_DIR, "08_treatment_per_cluster.png"))

# ── Slide 13: Summary ─────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
title_bar(slide, "Summary & Next Steps")
div = slide.shapes.add_shape(1, Inches(6.7), Inches(1.3), Inches(0.03), Inches(5.5))
div.fill.solid(); div.fill.fore_color.rgb = RGBColor(0xCC,0xDD,0xEE); div.line.fill.background()

def body_text(slide, left, top, width, height, lines, font_size=13):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run(); r.text = line
        if line.startswith("\u2022") or line.startswith("  \u2013"):
            r.font.size = Pt(font_size); r.font.color.rgb = BODY_COLOR
        else:
            r.font.size = Pt(font_size); r.font.bold = True; r.font.color.rgb = TITLE_COLOR

body_text(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.5), [
    "Update 2 — key changes",
    "\u2022 ER+/ER- from patient metadata (Table S1)",
    "\u2022 Marker genes: highlight_genes.txt",
    f"  \u2013 {GENE_LABEL}",
    "\u2022 DEG: FindAllMarkers per cluster",
    "\u2022 GSEA: Hallmark + Oncogenic (C6)",
    "  \u2013 HVGs only for speed",
    "\u2022 Metadata: gender / bone site / treatment",
    "  \u2013 Cluster 11 = highlight-gene-high",
], font_size=13)

body_text(slide, Inches(7.0), Inches(1.3), Inches(6.0), Inches(5.5), [
    "Next steps",
    "\u2022 Annotate clusters with canonical markers",
    "  \u2013 Cross-ref Yuliang\u2019s signature",
    "\u2022 Link highlight-gene-high cluster",
    "  \u2013 to S2C2 crosstalk model",
    "\u2022 Validate in mouse scRNA (Jianting W2/W3)",
    "\u2022 Drug target screen on top cluster",
], font_size=13)

prs.save(OUT_PPTX)
print(f"Saved: {OUT_PPTX}  ({prs.slides.__len__()} slides)")
