"""
BM project-human data_update_3.pptx — 1 figure per slide
Figures from bm/figures/update3/
Key change: resolution=0.3 yields Cluster 1 = 94.8% BC_ER+ (1791/1889 cells)
Gene set: highlight_genes.txt UNION ESR1/SP1
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os, csv

FIG_DIR  = "bm/figures/update3"
OUT_PPTX = "BM project-human data_update_3.pptx"

# ── Gene labels ───────────────────────────────────────────────────────────────
_GENE_FILE = "bm/figures/update1/highlight_genes.txt"
try:
    raw = open(_GENE_FILE).read().strip()
    HIGHLIGHT_GENES = [g.strip() for g in raw.split(",") if g.strip()]
except FileNotFoundError:
    HIGHLIGHT_GENES = ["Timp2", "Mmp14", "Pdgfa", "Fn1", "Vim"]
HL_LABEL   = " / ".join(HIGHLIGHT_GENES)
ES_LABEL   = "ESR1 / SP1"
COMB_LABEL = HL_LABEL + " / " + ES_LABEL

# ── Best params (from 00_grid_scores.csv) ─────────────────────────────────────
BEST_PARAMS = "n.neighbors=10  min.dist=0.05  resolution=0.3  n_clusters=11"
CL1_FINDING = "Cluster 1: 1791/1889 cells = 94.8% BC_ER+"

TITLE_COLOR = RGBColor(0x1F, 0x49, 0x7D)
BODY_COLOR  = RGBColor(0x26, 0x26, 0x26)
GRAY        = RGBColor(0x88, 0x88, 0x88)
ACCENT      = RGBColor(0xD6, 0x27, 0x28)   # red for highlight callout

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = TITLE_COLOR; bar.line.fill.background()
    txb = slide.shapes.add_textbox(Inches(0.3), Inches(0.12), Inches(12.0), Inches(0.55))
    p = txb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
    if subtitle:
        txb2 = slide.shapes.add_textbox(Inches(0.3), Inches(0.68), Inches(12.0), Inches(0.35))
        p2 = txb2.text_frame.paragraphs[0]
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(13); r2.font.color.rgb = RGBColor(0xCC,0xDD,0xEE)

def add_img(slide, path, left, top, width, height):
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)
    else:
        box = slide.shapes.add_shape(1, left, top, width, height)
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF0,0xF4,0xF8)
        box.line.color.rgb = RGBColor(0xAA,0xBB,0xCC)
        txb = slide.shapes.add_textbox(left+Inches(0.1), top+Inches(0.1),
                                        width-Inches(0.2), height-Inches(0.2))
        r = txb.text_frame.paragraphs[0].add_run()
        r.text = f"[Missing: {os.path.basename(path)}]"
        r.font.size = Pt(11); r.font.color.rgb = GRAY

def fig_slide(title, img_path):
    slide = prs.slides.add_slide(blank)
    title_bar(slide, title)
    add_img(slide, img_path, Inches(0.4), Inches(1.2), Inches(12.5), Inches(6.1))
    return slide

def callout_box(slide, text, left, top, width, height, color=ACCENT):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFF,0xF0,0xF0)
    box.line.color.rgb = color
    txb = slide.shapes.add_textbox(left+Inches(0.08), top+Inches(0.05),
                                    width-Inches(0.16), height-Inches(0.1))
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = color

# ── Slide 1: Title ────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
bg = slide.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(7.5))
bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0xF7,0xF9,0xFC); bg.line.fill.background()
bar = slide.shapes.add_shape(1, 0, Inches(2.8), Inches(13.33), Inches(2.0))
bar.fill.solid(); bar.fill.fore_color.rgb = TITLE_COLOR; bar.line.fill.background()
txb = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.0), Inches(1.2))
p = txb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "BM Project \u2013 Human Data Update 3"
r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
txb2 = slide.shapes.add_textbox(Inches(1.0), Inches(4.1), Inches(11.0), Inches(0.5))
p2 = txb2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "Parameter sweep clustering | BC_ER+ separation | Highlight + ESR1/SP1 analysis"
r2.font.size = Pt(15); r2.font.color.rgb = RGBColor(0xCC,0xDD,0xEE)
txb3 = slide.shapes.add_textbox(Inches(1.0), Inches(6.5), Inches(11.0), Inches(0.5))
p3 = txb3.text_frame.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run(); r3.text = "Jianting Sheng | Houston Methodist Neal Cancer Center"
r3.font.size = Pt(13); r3.font.color.rgb = GRAY

# ── Slide 2: Sanity check — cells per cancer type ────────────────────────────
slide = prs.slides.add_slide(blank)
title_bar(slide, "Sanity Check \u2014 Macrophage Cells per Cancer Type")
add_img(slide, os.path.join(FIG_DIR, "00_sanity_cells_per_cancer.png"),
        Inches(1.5), Inches(1.2), Inches(10.3), Inches(6.1))

# ── Slide 3: Sanity check — cells per sample ─────────────────────────────────
slide = prs.slides.add_slide(blank)
title_bar(slide, "Sanity Check \u2014 Macrophage Cells per Sample")
add_img(slide, os.path.join(FIG_DIR, "00_sanity_cells_per_sample.png"),
        Inches(0.3), Inches(1.2), Inches(12.7), Inches(6.1))

# ── Slide 4: UMAP overview ────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
title_bar(slide, "Macrophage Sub-clustering UMAP \u2014 Best Params")
add_img(slide, os.path.join(FIG_DIR, "01_umap_overview.png"),
        Inches(0.4), Inches(1.2), Inches(12.5), Inches(6.1))

fig_slide("BC / BC_ER Cells Highlighted on UMAP",
          os.path.join(FIG_DIR, "04_bc_highlight_umap.png"))

fig_slide("Cancer Type Composition per Macrophage Cluster",
          os.path.join(FIG_DIR, "02_cancer_type_bar_bc_er.png"))

fig_slide(f"Feature Plots \u2014 Highlight Genes",
          os.path.join(FIG_DIR, "03a_feature_plots_highlight.png"))

fig_slide(f"Feature Plots \u2014 ESR1 / SP1",
          os.path.join(FIG_DIR, "03b_feature_plots_esr1_sp1.png"))

fig_slide(f"Combined Dot Plot \u2014 Highlight + ESR1/SP1",
          os.path.join(FIG_DIR, "03c_dotplot_combined.png"))

fig_slide("DEG Dot Plot \u2014 Top 3 Marker Genes per Cluster",
          os.path.join(FIG_DIR, "04b_deg_dotplot.png"))

fig_slide(f"Highlight + ESR1/SP1 Expression \u2014 Top 6 Clusters",
          os.path.join(FIG_DIR, "04c_dotplot_comb_top_clusters.png"))

fig_slide("GSEA \u2014 All Databases | Top 5 Up/Down per Cluster",
          os.path.join(FIG_DIR, "05e_gsea_alldb_stacking.png"))

fig_slide("GSEA \u2014 Top Pathways per Cluster (Bubble Plot)",
          os.path.join(FIG_DIR, "05b_gsea_bubble.png"))

fig_slide("GSEA (Broader) \u2014 Cluster 1 (BC_ER+)",
          os.path.join(FIG_DIR, "05f_cl1_combined_nofacet.png"))

fig_slide("Bone Site / Tissue Origin per Macrophage Sub-cluster",
          os.path.join(FIG_DIR, "07_tissue_origin_per_cluster.png"))

fig_slide("Treatment Response per Macrophage Sub-cluster",
          os.path.join(FIG_DIR, "08_treatment_per_cluster.png"))

# ── Slide 17: Summary ─────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
title_bar(slide, "Summary \u2014 Update 3")
div = slide.shapes.add_shape(1, Inches(6.7), Inches(1.3), Inches(0.03), Inches(5.5))
div.fill.solid(); div.fill.fore_color.rgb = RGBColor(0xCC,0xDD,0xEE); div.line.fill.background()

def body_text(slide, left, top, width, height, lines, font_size=13):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run(); r.text = line
        if line.startswith("\u2022") or line.startswith("  \u2013"):
            r.font.size = Pt(font_size); r.font.color.rgb = BODY_COLOR
        else:
            r.font.size = Pt(font_size); r.font.bold = True; r.font.color.rgb = TITLE_COLOR

body_text(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.8), [
    "Clustering",
    f"\u2022 {CL1_FINDING}",
    "  \u2013 res=0.3 dominant; UMAP params irrelevant",
    "  \u2013 BC08 = 54% of all BC_ER+ cells (single-patient dominance)",
    "Gene expression",
    "\u2022 Cluster 10 highest for TIMP2/MMP14/PDGFA/FN1/VIM",
    "  \u2013 Cluster 1 (BC_ER+) ranks 3rd",
    "\u2022 ESR1/SP1 co-expression not exclusive to Cluster 1",
    "Hallmark / Oncogenic GSEA (H + C6)",
    "\u2022 Cluster 1: no significant terms (all padj>0.15)",
    "  \u2013 Trending down: IL-6/JAK/STAT3, IFN-g, KRAS-breast",
    "\u2022 Cluster 0: IFN-a/g strongly UP (padj<0.003)",
    "\u2022 Cluster 4 (LC): KRAS + mTORC1 + TNFa/NFkB UP",
    "\u2022 Cluster 8: STK33 DN, Allograft rejection UP",
    "\u2022 Cluster 9: EMT UP (only cluster)",
], font_size=10)

body_text(slide, Inches(7.0), Inches(1.3), Inches(6.0), Inches(5.8), [
    "Broader GSEA (GO:BP + Reactome + C7-IMMUNESIGDB)",
    "\u2022 Cluster 1 (BC_ER+) — 21 deduplicated terms:",
    "  \u2013 Mac. polarization: PPARG-KO, WBP7-KO, MDSC DOWN",
    "    (all padj<0.05 — strongest signals in cluster 1)",
    "  \u2013 Cytokine: IFN-g signaling, cytokine network DOWN",
    "  \u2013 TGF-b: TH17/TGF-b/IL-6, AKT-TGF-b DOWN",
    "  \u2013 Immunosuppression: T cell activation DOWN",
    "  \u2013 Wnt/osteogenic: Wnt signaling (Reactome) DOWN",
    "  \u2013 ECM: ECM organization weakly UP",
    "Interpretation re: BM manuscript",
    "\u2022 BC_ER+ macrophages appear transcriptionally suppressed",
    "  \u2013 Consistent with immunosuppressive ER+ TME",
    "  \u2013 Down-regulated Wnt may reflect niche dependency",
    "  \u2013 PPARG/WBP7/MDSC signatures: non-classical activation",
    "Next steps",
    "\u2022 DEG: Cluster 1 vs BC_ER- clusters",
    "\u2022 Validate SP1/ESR1 co-expression in Cluster 1",
    "\u2022 Link to S2C2 crosstalk; validate in mouse W2/W3",
], font_size=10)

callout_box(slide, f"\u2605 {CL1_FINDING}  |  BC_ER+ macrophages: immunosuppressed, non-classical, Wnt/cytokine DOWN",
            Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.45))

prs.save(OUT_PPTX)
print(f"Saved: {OUT_PPTX}  ({len(prs.slides)} slides)")
