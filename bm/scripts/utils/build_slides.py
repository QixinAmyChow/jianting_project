"""
BM project-human data_update_3.pptx
Reads figures from bm_analysis_out/figures_update_3/ and assembles the presentation.
Run after pipeline_update3.R completes.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os, csv

FIG_DIR  = "bm_analysis_out/figures_update_3"
OUT_PPTX = "BM project-human data_update_3.pptx"

BEST_PARAMS = "n.neighbors=10  min.dist=0.05  resolution=0.3  n_clusters=11"
CL1_FINDING = "Cluster 1: 1791/1889 cells = 94.8% BC_ER+"

TITLE_COLOR = RGBColor(0x1F, 0x49, 0x7D)
BODY_COLOR  = RGBColor(0x26, 0x26, 0x26)
GRAY        = RGBColor(0x88, 0x88, 0x88)
ACCENT      = RGBColor(0xD6, 0x27, 0x28)

DB_LEGEND   = "HALLMARK (blue) · C6 (lt.blue) · GO:BP (green) · Reactome (orange) · C7-IMMUNESIGDB (purple)"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = TITLE_COLOR; bar.line.fill.background()
    txb = slide.shapes.add_textbox(Inches(0.3), Inches(0.12), Inches(12.0), Inches(0.55))
    p = txb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    if subtitle:
        txb2 = slide.shapes.add_textbox(Inches(0.3), Inches(0.68), Inches(12.0), Inches(0.35))
        p2 = txb2.text_frame.paragraphs[0]
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(13); r2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)


def add_img(slide, path, left, top, width, height):
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)
    else:
        box = slide.shapes.add_shape(1, left, top, width, height)
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
        box.line.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
        r = box.text_frame.paragraphs[0].add_run()  # note: shapes don't have text_frame; use textbox
        txb = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.1),
                                       width - Inches(0.2), height - Inches(0.2))
        r2 = txb.text_frame.paragraphs[0].add_run()
        r2.text = f"[Missing: {os.path.basename(path)}]"
        r2.font.size = Pt(11); r2.font.color.rgb = GRAY


def fig_slide(title, img_path, subtitle=None):
    slide = prs.slides.add_slide(blank)
    title_bar(slide, title, subtitle)
    add_img(slide, img_path, Inches(0.4), Inches(1.2), Inches(12.5), Inches(6.1))
    return slide


def callout(slide, text, left, top, width, height):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xF0)
    box.line.color.rgb = ACCENT
    txb = slide.shapes.add_textbox(left + Inches(0.08), top + Inches(0.05),
                                   width - Inches(0.16), height - Inches(0.1))
    tf = txb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = text
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = ACCENT


def body_text(slide, left, top, width, height, lines, font_size=13):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        r = p.add_run(); r.text = line
        if line.startswith("•") or line.startswith("  –"):
            r.font.size = Pt(font_size); r.font.color.rgb = BODY_COLOR
        else:
            r.font.size = Pt(font_size); r.font.bold = True; r.font.color.rgb = TITLE_COLOR


# ── Slide 1: Title ────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
bg = slide.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(7.5))
bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFC); bg.line.fill.background()
bar = slide.shapes.add_shape(1, 0, Inches(2.8), Inches(13.33), Inches(2.0))
bar.fill.solid(); bar.fill.fore_color.rgb = TITLE_COLOR; bar.line.fill.background()
txb = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.0), Inches(1.2))
p = txb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "BM Project – Human Data Update 3"
r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
txb2 = slide.shapes.add_textbox(Inches(1.0), Inches(4.1), Inches(11.0), Inches(0.5))
p2 = txb2.text_frame.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "Parameter sweep clustering | BC_ER+ separation | Highlight + ESR1/SP1 analysis"
r2.font.size = Pt(15); r2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
txb3 = slide.shapes.add_textbox(Inches(1.0), Inches(6.5), Inches(11.0), Inches(0.5))
p3 = txb3.text_frame.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run(); r3.text = "Jianting Sheng | Houston Methodist Neal Cancer Center"
r3.font.size = Pt(13); r3.font.color.rgb = GRAY

# ── Slides 2–3: Sanity checks ─────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
title_bar(slide, "Sanity Check — Macrophage Cells per Cancer Type")
add_img(slide, os.path.join(FIG_DIR, "00_sanity_cells_per_cancer.png"),
        Inches(1.5), Inches(1.2), Inches(10.3), Inches(6.1))

slide = prs.slides.add_slide(blank)
title_bar(slide, "Sanity Check — Macrophage Cells per Sample")
add_img(slide, os.path.join(FIG_DIR, "00_sanity_cells_per_sample.png"),
        Inches(0.3), Inches(1.2), Inches(12.7), Inches(6.1))

# ── Slides 4–5: UMAP + BC highlight ──────────────────────────────────────────
slide = prs.slides.add_slide(blank)
title_bar(slide, "Macrophage Sub-clustering UMAP — Best Params", BEST_PARAMS)
add_img(slide, os.path.join(FIG_DIR, "01_umap_overview.png"),
        Inches(0.4), Inches(1.2), Inches(12.5), Inches(6.1))

fig_slide("BC / BC_ER Cells Highlighted on UMAP",
          os.path.join(FIG_DIR, "04_bc_highlight_umap.png"))

# ── Slide 6: Cancer type bar ──────────────────────────────────────────────────
fig_slide("Cancer Type Composition per Macrophage Cluster",
          os.path.join(FIG_DIR, "02_cancer_type_bar_bc_er.png"))

# ── Slides 7–9: Gene expression ───────────────────────────────────────────────
fig_slide("Feature Plots — Highlight Genes",
          os.path.join(FIG_DIR, "03a_feature_plots_highlight.png"))
fig_slide("Feature Plots — ESR1 / SP1",
          os.path.join(FIG_DIR, "03b_feature_plots_esr1_sp1.png"))
fig_slide("Combined Dot Plot — Highlight + ESR1/SP1",
          os.path.join(FIG_DIR, "03c_dotplot_combined.png"))

# ── Slides 10–11: DEG ─────────────────────────────────────────────────────────
fig_slide("DEG Dot Plot — Top 3 Marker Genes per Cluster",
          os.path.join(FIG_DIR, "04b_deg_dotplot.png"))
fig_slide("Highlight + ESR1/SP1 Expression — Top 6 Clusters",
          os.path.join(FIG_DIR, "04c_dotplot_comb_top_clusters.png"))

# ── Slides 12–13: GSEA all clusters ──────────────────────────────────────────
fig_slide("GSEA — All Databases | Top 5 Up/Down per Cluster",
          os.path.join(FIG_DIR, "05e_gsea_alldb_stacking.png"), DB_LEGEND)
fig_slide("GSEA — Top Pathways per Cluster (Bubble Plot)",
          os.path.join(FIG_DIR, "05b_gsea_bubble.png"))

# ── Slides 14–15: Cluster 1 GSEA ─────────────────────────────────────────────
fig_slide("GSEA (Broader) — Cluster 1 (BC_ER+) by Hypothesis",
          os.path.join(FIG_DIR, "05f_cl1_combined_nofacet.png"))
fig_slide("GSEA — Cluster 1 Hypothesis Groups (Colored by Direction)",
          os.path.join(FIG_DIR, "05c_cluster1_hypo_gsea.png"))

# ── Slides 16–18: Clinical metadata ──────────────────────────────────────────
fig_slide("Bone Site / Tissue Origin per Macrophage Sub-cluster",
          os.path.join(FIG_DIR, "07_tissue_origin_per_cluster.png"))
fig_slide("Treatment Response per Macrophage Sub-cluster",
          os.path.join(FIG_DIR, "08_treatment_per_cluster.png"))
fig_slide("Gender Distribution per Macrophage Sub-cluster",
          os.path.join(FIG_DIR, "06_gender_per_cluster.png"))

# ── Slide 19: Summary ─────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank)
title_bar(slide, "Summary — Update 3")
div = slide.shapes.add_shape(1, Inches(6.7), Inches(1.3), Inches(0.03), Inches(5.5))
div.fill.solid(); div.fill.fore_color.rgb = RGBColor(0xCC, 0xDD, 0xEE); div.line.fill.background()

body_text(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.8), [
    "Clustering",
    f"• {CL1_FINDING}",
    "  – res=0.3 dominant; UMAP params irrelevant",
    "  – BC08 = 54% of all BC_ER+ cells (single-patient dominance)",
    "Gene expression",
    "• Cluster 10 highest for TIMP2/MMP14/PDGFA/FN1/VIM",
    "  – Cluster 1 (BC_ER+) ranks 3rd",
    "• ESR1/SP1 co-expression not exclusive to Cluster 1",
    "Hallmark / Oncogenic GSEA (H + C6)",
    "• Cluster 1: no significant terms (all padj>0.15)",
    "  – Trending down: IL-6/JAK/STAT3, IFN-g, KRAS-breast",
    "• Cluster 0: IFN-a/g strongly UP (padj<0.003)",
    "• Cluster 4 (LC): KRAS + mTORC1 + TNFa/NFkB UP",
    "• Cluster 8: STK33 DN, Allograft rejection UP",
    "• Cluster 9: EMT UP (only cluster)",
], font_size=10)

body_text(slide, Inches(7.0), Inches(1.3), Inches(6.0), Inches(5.8), [
    "Broader GSEA (GO:BP + Reactome + C7-IMMUNESIGDB)",
    "• Cluster 1 (BC_ER+) — 21 deduplicated terms:",
    "  – Mac. polarization: PPARG-KO, WBP7-KO, MDSC DOWN",
    "    (all padj<0.05 — strongest signals in cluster 1)",
    "  – Cytokine: IFN-g signaling, cytokine network DOWN",
    "  – TGF-b: TH17/TGF-b/IL-6, AKT-TGF-b DOWN",
    "  – Immunosuppression: T cell activation DOWN",
    "  – Wnt/osteogenic: Wnt signaling (Reactome) DOWN",
    "  – ECM: ECM organization weakly UP",
    "Interpretation re: BM manuscript",
    "• BC_ER+ macrophages appear transcriptionally suppressed",
    "  – Consistent with immunosuppressive ER+ TME",
    "  – Down-regulated Wnt may reflect niche dependency",
    "  – PPARG/WBP7/MDSC signatures: non-classical activation",
    "Next steps",
    "• DEG: Cluster 1 vs BC_ER- clusters",
    "• Validate SP1/ESR1 co-expression in Cluster 1",
    "• Link to S2C2 crosstalk; validate in mouse W2/W3",
], font_size=10)

callout(slide,
        f"★ {CL1_FINDING}  |  BC_ER+ macrophages: immunosuppressed, non-classical, Wnt/cytokine DOWN",
        Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.45))

prs.save(OUT_PPTX)
print(f"Saved: {OUT_PPTX}  ({len(prs.slides)} slides)")
