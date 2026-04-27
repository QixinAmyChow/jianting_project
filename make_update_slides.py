from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

TITLE_COLOR = RGBColor(0x1F, 0x49, 0x7D)  # dark blue
BODY_COLOR  = RGBColor(0x26, 0x26, 0x26)
GRAY        = RGBColor(0x88, 0x88, 0x88)
ACCENT      = RGBColor(0xC0, 0x39, 0x2B)  # red accent

blank_layout = prs.slide_layouts[6]  # completely blank

def add_title_bar(slide, title_text, subtitle_text=None):
    # colored bar across top
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TITLE_COLOR
    bar.line.fill.background()

    txb = slide.shapes.add_textbox(Inches(0.3), Inches(0.12), Inches(12.0), Inches(0.55))
    tf = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    if subtitle_text:
        txb2 = slide.shapes.add_textbox(Inches(0.3), Inches(0.68), Inches(12.0), Inches(0.35))
        tf2 = txb2.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle_text
        run2.font.size = Pt(13)
        run2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)

def add_body_text(slide, left, top, width, height, lines, font_size=13):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        if line.startswith("•"):
            p.level = 1
            run = p.add_run()
            run.text = line
            run.font.size = Pt(font_size)
            run.font.color.rgb = BODY_COLOR
        elif line.startswith("  –"):
            p.level = 2
            run = p.add_run()
            run.text = line.strip()
            run.font.size = Pt(font_size - 1)
            run.font.color.rgb = GRAY
        else:
            run = p.add_run()
            run.text = line
            run.font.size = Pt(font_size)
            run.font.bold = True
            run.font.color.rgb = TITLE_COLOR

def add_placeholder_box(slide, left, top, width, height, label, note=None):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
    box.line.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)

    txb = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.15), width - Inches(0.2), height - Inches(0.3))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = GRAY

    if note:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = note
        r2.font.size = Pt(10)
        r2.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)

# ── Slide 1: Title ───────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
bg = slide.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(7.5))
bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFC); bg.line.fill.background()

bar = slide.shapes.add_shape(1, 0, Inches(2.8), Inches(13.33), Inches(2.0))
bar.fill.solid(); bar.fill.fore_color.rgb = TITLE_COLOR; bar.line.fill.background()

txb = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.0), Inches(1.2))
tf = txb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "BM Project – Human Data Update"
run.font.size = Pt(32); run.font.bold = True; run.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)

txb2 = slide.shapes.add_textbox(Inches(1.0), Inches(4.1), Inches(11.0), Inches(0.5))
tf2 = txb2.text_frame
p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "Macrophage Sub-clustering Analysis | GSE266330"
r2.font.size = Pt(16); r2.font.color.rgb = RGBColor(0xCC,0xDD,0xEE)

txb3 = slide.shapes.add_textbox(Inches(1.0), Inches(6.5), Inches(11.0), Inches(0.5))
tf3 = txb3.text_frame
p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
r3 = p3.add_run()
r3.text = "Jianting Sheng | Houston Methodist Neal Cancer Center"
r3.font.size = Pt(13); r3.font.color.rgb = GRAY

# ── Slide 2: Analysis overview ───────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Analysis Overview", "Goal: Sub-cluster all macrophages; identify ER+/SP1-high clusters")

add_body_text(slide, Inches(0.4), Inches(1.3), Inches(5.8), Inches(5.5), [
    "Dataset",
    "• Liu et al. 2025 (Cell Genomics)",
    "  – 42 human BM samples, 8 cancer types",
    "  – GEO: GSE266330",
    "  – ~180K high-quality cells",
    "",
    "Cell selection",
    "• Pool macrophages from all samples",
    "• Re-cluster (standard HVG pipeline)",
    "• Identify sub-populations",
])

add_body_text(slide, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.5), [
    "Questions to answer",
    "• Which cluster has ER + SP1 both high?",
    "  – Match Yuliang's gene signature",
    "• What cancer types dominate each cluster?",
    "  – BC vs. TC vs. other origin %",
    "• Is any cluster enriched in metastasis",
    "  – vs. control?",
    "• Is treatment response associated with",
    "  – a specific macrophage sub-population?",
    "• Gender and tissue origin breakdown",
])

# divider
div = slide.shapes.add_shape(1, Inches(6.5), Inches(1.3), Inches(0.03), Inches(5.5))
div.fill.solid(); div.fill.fore_color.rgb = RGBColor(0xCC,0xDD,0xEE); div.line.fill.background()

# ── Slide 3: Macrophage sub-clustering UMAP ──────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Macrophage Sub-clustering (UMAP)", "All samples pooled → re-clustered on high-variance genes")

add_placeholder_box(slide, Inches(0.4), Inches(1.3), Inches(5.5), Inches(5.5),
    "[UMAP — macrophage sub-clusters]",
    "Colored by cluster ID\nLabel: cluster 0–N")

add_placeholder_box(slide, Inches(6.2), Inches(1.3), Inches(6.7), Inches(2.5),
    "[UMAP — colored by cancer type of origin]",
    "BC / TC / Lung / Kidney / Other")

add_placeholder_box(slide, Inches(6.2), Inches(4.0), Inches(6.7), Inches(2.8),
    "[Stacked bar — % cancer type per cluster]",
    "X: cluster; Y: fraction; color: cancer type")

# ── Slide 4: ER + SP1 expression ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "ER / SP1 Co-expression in Macrophage Sub-clusters",
              "Target: identify cluster(s) where both ESR1 (ER) and SP1 are high")

add_placeholder_box(slide, Inches(0.4), Inches(1.3), Inches(4.0), Inches(5.5),
    "[Feature plot — ESR1 expression]",
    "UMAP; color = log-normalized counts")

add_placeholder_box(slide, Inches(4.6), Inches(1.3), Inches(4.0), Inches(5.5),
    "[Feature plot — SP1 expression]",
    "UMAP; color = log-normalized counts")

add_placeholder_box(slide, Inches(8.9), Inches(1.3), Inches(4.0), Inches(5.5),
    "[Dot plot — ESR1 & SP1 across clusters]",
    "Dot size = % expressing; color = avg expr")

# ── Slide 5: Metastasis vs. Control in macrophage clusters ───────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Macrophage Cluster Enrichment — Metastasis vs. Control",
              "Are specific sub-clusters over-represented in metastatic samples?")

add_placeholder_box(slide, Inches(0.4), Inches(1.3), Inches(6.2), Inches(5.5),
    "[Bar / violin — cluster proportion per condition]",
    "Metastasis vs. control; stats overlay")

add_placeholder_box(slide, Inches(6.9), Inches(1.3), Inches(6.0), Inches(5.5),
    "[Heatmap — top DEGs per cluster, met vs ctrl]",
    "Rows: genes; columns: clusters\nHighlight ER+ / SP1-high cluster")

# ── Slide 6: Treatment response association ───────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Macrophage Sub-population & Treatment Response",
              "Is any macrophage sub-cluster associated with treatment response?")

add_placeholder_box(slide, Inches(0.4), Inches(1.3), Inches(6.2), Inches(5.5),
    "[UMAP / bar — colored by treatment response]",
    "Responder vs. non-responder\n(if metadata available in GSE266330)")

add_body_text(slide, Inches(6.9), Inches(1.3), Inches(6.0), Inches(2.5), [
    "Metadata to extract from GSE266330",
    "• Treatment history (chemo / hormone / immuno)",
    "• Response classification (CR/PR/SD/PD)",
    "• Time to progression",
], font_size=12)

add_body_text(slide, Inches(6.9), Inches(3.9), Inches(6.0), Inches(2.8), [
    "Analysis plan",
    "• Correlate cluster % with response label",
    "• Survival analysis if follow-up data exists",
    "• Focus on ER+/SP1-high cluster",
], font_size=12)

# ── Slide 7: Cancer origin & clinical metadata ────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Cancer Origin, Gender & Tissue Site per Macrophage Cluster",
              "细分之后是否有更细的富集的cancer type / gender / tissue origin")

add_placeholder_box(slide, Inches(0.4), Inches(1.3), Inches(4.0), Inches(5.5),
    "[Pie / stacked bar — BC vs TC vs other\nper macrophage sub-cluster]",
    "Highlight dominant origin per cluster")

add_placeholder_box(slide, Inches(4.6), Inches(1.3), Inches(4.0), Inches(5.5),
    "[Bar — gender distribution per cluster]",
    "Male / Female fraction")

add_placeholder_box(slide, Inches(8.9), Inches(1.3), Inches(4.0), Inches(5.5),
    "[Bar — tissue/site of metastasis per cluster]",
    "Bone niche location if available")

# ── Slide 8: Summary & next steps ────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
add_title_bar(slide, "Summary & Next Steps")

add_body_text(slide, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.5), [
    "Expected findings",
    "• N macrophage sub-clusters identified",
    "  – At least 1 cluster ER+/SP1-high",
    "• ER+/SP1-high cluster likely enriched in:",
    "  – Breast cancer (ER+ subtype)",
    "  – Metastatic vs. control condition",
    "• Mφ-OC ecosystem (Liu 2025) maps to",
    "  – specific sub-cluster(s) here",
], font_size=13)

add_body_text(slide, Inches(7.0), Inches(1.3), Inches(6.0), Inches(5.5), [
    "Next steps",
    "• Annotate clusters with marker genes",
    "  – Cross-ref Yuliang's signature",
    "• Link to S2C2 crosstalk model",
    "  – ER+ Mφ → CD8+ T cell axis",
    "• Validate in mouse scRNA (Jianting W2/W3)",
    "• Drug target screen on ER+/SP1-high cluster",
], font_size=13)

div2 = slide.shapes.add_shape(1, Inches(6.7), Inches(1.3), Inches(0.03), Inches(5.5))
div2.fill.solid(); div2.fill.fore_color.rgb = RGBColor(0xCC,0xDD,0xEE); div2.line.fill.background()

prs.save("BM project-human data_update.pptx")
print("Saved: BM project-human data_update.pptx")
